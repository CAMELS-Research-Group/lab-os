#!/usr/bin/env python
"""Build the Building-With-Claude kickoff deck from its markdown source.

Single source of truth: `docs/workshops/building/kickoff-deck.md`. This script
parses that file and renders `building-with-claude-kickoff.pptx` — so the deck is
edited in the markdown, never in the binary. The `.pptx`/`.odp` renders are
gitignored; regenerate them with this script.

Usage:
    python tools/build_deck.py                 # src + out at their defaults
    python tools/build_deck.py --out deck.pptx # custom output path
    python tools/build_deck.py path/to/deck.md # custom source

Export to ODP (optional, needs LibreOffice):
    soffice --headless --convert-to odp --outdir <dir> <the .pptx>

The parser is deliberately strict: it asserts the source has exactly the
expected slide count and contiguous S-numbers, and it prints an overflow
estimate per content slide (a `TIGHT` warning means the body may not fit).
A markdown edit that breaks the structure fails the build loudly rather than
producing a silently-wrong deck.
"""
import argparse
import math
import re
import sys
from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

REPO = Path(__file__).resolve().parent.parent
DEFAULT_SRC = REPO / "docs" / "workshops" / "building" / "kickoff-deck.md"
DEFAULT_OUT = REPO / "docs" / "workshops" / "building" / "building-with-claude-kickoff.pptx"
EXPECTED_SLIDES = 27

# ---- lab palette ---------------------------------------------------------
INDIGO    = RGBColor(0x5B, 0x54, 0xE8)
INDIGO_DK = RGBColor(0x25, 0x1D, 0xD3)
INK       = RGBColor(0x1A, 0x1A, 0x2E)
GREY      = RGBColor(0x6B, 0x70, 0x80)
WHITE     = RGBColor(0xFF, 0xFF, 0xFF)
MIST      = RGBColor(0xEC, 0xEB, 0xFC)
FONT      = "Calibri"


# ======================== markdown parsing ================================

def _strip_inline(text):
    """Drop inline code backticks and emphasis markers; this renderer has no
    inline rich-text, so `code`, **bold**, and *italic* become plain text.
    Bold is unwrapped before italic so `**x**` doesn't leave a stray `*`."""
    text = text.replace("`", "")
    text = re.sub(r"\*\*(.+?)\*\*", r"\1", text)
    text = re.sub(r"\*(.+?)\*", r"\1", text)
    return text.strip()


def _classify(item):
    """Return (kind, clean_text) for a list item's raw text.
    A whole-item **bold** is a caption (the payoff line); a whole-item *italic*
    is the lede; anything else is a plain bullet."""
    m = re.match(r"^\*\*(.+)\*\*$", item)
    if m:
        return "caption", _strip_inline(m.group(1))
    m = re.match(r"^\*(.+)\*$", item)
    if m:
        return "lede", _strip_inline(m.group(1))
    return "bullet", _strip_inline(item)


def _split_blocks(md):
    """Yield (num, title_raw, body_lines) per `### S<n> · <title>` slide block."""
    lines = md.splitlines()
    header = re.compile(r"^### S(\d+) · (.+?)\s*$")
    blocks = []
    cur = None
    for ln in lines:
        m = header.match(ln)
        if m:
            if cur:
                blocks.append(cur)
            cur = (int(m.group(1)), m.group(2).strip(), [])
        elif cur is not None:
            if ln.startswith("## ") or ln.startswith("### "):
                blocks.append(cur)
                cur = None
            else:
                cur[2].append(ln)
    if cur:
        blocks.append(cur)
    return blocks


def _extract_onslide_notes(body_lines):
    """Split a slide body into (onslide_inline, onslide_item_lines, notes_text)."""
    mode = None
    inline = ""
    onslide_lines = []
    notes_parts = []
    for ln in body_lines:
        s = ln.strip()
        if s.startswith("**On-slide"):
            mode = "onslide"
            inline = re.sub(r"^\*\*On-slide.*?\*\*", "", s).strip()
            continue
        if s.startswith("**Notes"):
            mode = "notes"
            first = re.sub(r"^\*\*Notes.*?\*\*", "", s).strip().lstrip(":").strip()
            if first:
                notes_parts.append(first)
            continue
        if mode == "onslide":
            if s and s != "---":
                onslide_lines.append(ln)
        elif mode == "notes":
            if s == "" or s == "---":
                mode = "done"
            else:
                notes_parts.append(s)
    notes = _strip_inline(" ".join(notes_parts)) if notes_parts else ""
    # _strip_inline collapses backticks/bold; re-collapse whitespace it left
    notes = re.sub(r"\s+", " ", notes).strip()
    return inline, onslide_lines, notes


def _parse_items(onslide_lines):
    """Parse list/numbered lines into [(text, numbered)], joining continuations."""
    items = []
    for ln in onslide_lines:
        s = ln.strip()
        if not s:
            continue
        mnum = re.match(r"^(\d+)\.\s+(.*)$", s)
        if s.startswith("- "):
            items.append([s[2:].strip(), False])
        elif mnum:
            items.append([mnum.group(2).strip(), True])
        elif items:  # wrapped continuation of the previous item
            items[-1][0] += " " + s
    return items


def parse_deck(md):
    """Parse the kickoff-deck markdown into an ordered list of slide dicts."""
    slides = []
    for num, title_raw, body in _split_blocks(md):
        inline, onslide_lines, notes = _extract_onslide_notes(body)

        if title_raw == "Title":
            items = _parse_items(onslide_lines)
            title, org, subs = "", "CAMELS Research Group", []
            for text, _ in items:
                kind, clean = _classify(text)
                if kind == "caption":
                    title = clean
                elif kind == "lede":
                    org = clean
                else:
                    subs.append(clean)
            slides.append({"num": num, "kind": "title", "title": title,
                           "subs": subs, "org": org, "notes": notes})

        elif title_raw.startswith("Divider"):
            full = (inline + " " + " ".join(l.strip() for l in onslide_lines)).strip()
            left, _, sub_part = full.partition(" — ")
            kicker, _, dtitle = left.partition(" · ")
            sub = _strip_inline(sub_part)
            slides.append({"num": num, "kind": "divider", "kicker": kicker.strip(),
                           "title": dtitle.strip(), "sub": sub, "notes": notes})

        else:
            items = _parse_items(onslide_lines)
            lede, bullets = "", []
            numbered = any(n for _, n in items)
            for i, (text, _) in enumerate(items):
                kind, clean = _classify(text)
                if i == 0 and kind == "lede":
                    lede = clean
                elif kind == "caption":
                    bullets.append("~" + clean)
                else:
                    bullets.append(clean)
            slides.append({"num": num, "kind": "content", "title": title_raw,
                           "lede": lede, "bullets": bullets, "notes": notes,
                           "marker": "" if numbered else "•   "})
    return slides


# ======================== pptx rendering ==================================

class Deck:
    def __init__(self):
        self.prs = Presentation()
        self.prs.slide_width = Inches(13.333)
        self.prs.slide_height = Inches(7.5)
        self.SW, self.SH = self.prs.slide_width, self.prs.slide_height
        self.blank = self.prs.slide_layouts[6]
        self.warnings = []

    def _rect(self, slide, l, t, w, h, color):
        sp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, l, t, w, h)
        sp.fill.solid(); sp.fill.fore_color.rgb = color
        sp.line.fill.background(); sp.shadow.inherit = False
        return sp

    def _tbox(self, slide, l, t, w, h, anchor=MSO_ANCHOR.TOP):
        tb = slide.shapes.add_textbox(l, t, w, h); tf = tb.text_frame
        tf.word_wrap = True; tf.vertical_anchor = anchor
        tf.margin_left = 0; tf.margin_right = 0; tf.margin_top = 0; tf.margin_bottom = 0
        return tf

    def _para(self, tf, text, size, color, bold=False, italic=False, after=8,
              first=False, align=PP_ALIGN.LEFT):
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        p.alignment = align; p.space_after = Pt(after); p.space_before = Pt(0)
        r = p.add_run(); r.text = text
        f = r.font; f.size = Pt(size); f.bold = bold; f.italic = italic
        f.name = FONT; f.color.rgb = color
        return p

    def _notes(self, slide, text):
        if text:
            slide.notes_slide.notes_text_frame.text = text

    def title_slide(self, s):
        sl = self.prs.slides.add_slide(self.blank)
        self._rect(sl, 0, 0, self.SW, self.SH, INDIGO)
        self._rect(sl, Inches(0.9), Inches(2.35), Inches(1.6), Inches(0.10), WHITE)
        tf = self._tbox(sl, Inches(0.9), Inches(2.6), Inches(11.6), Inches(3.0))
        self._para(tf, s["title"], 54, WHITE, bold=True, after=16, first=True)
        for ln in s["subs"]:
            self._para(tf, ln, 20, MIST, after=6)
        tf2 = self._tbox(sl, Inches(0.9), Inches(6.6), Inches(11.5), Inches(0.5))
        self._para(tf2, s["org"], 14, MIST, first=True)
        self._notes(sl, s["notes"])

    def divider_slide(self, s):
        sl = self.prs.slides.add_slide(self.blank)
        self._rect(sl, 0, 0, self.SW, self.SH, INK)
        self._rect(sl, 0, 0, Inches(0.28), self.SH, INDIGO)
        tf = self._tbox(sl, Inches(1.1), Inches(2.7), Inches(11.0), Inches(2.3))
        self._para(tf, s["kicker"].upper(), 16, INDIGO, bold=True, after=10, first=True)
        self._para(tf, s["title"], 44, WHITE, bold=True, after=8)
        if s["sub"]:
            self._para(tf, s["sub"], 20, GREY)
        self._notes(sl, s["notes"])

    def content_slide(self, s):
        title, lede, bullets = s["title"], s["lede"], s["bullets"]
        marker, idx = s["marker"], s["num"]
        sl = self.prs.slides.add_slide(self.blank)
        self._rect(sl, 0, 0, Inches(0.28), self.SH, INDIGO)
        tt = self._tbox(sl, Inches(0.95), Inches(0.58), Inches(11.6), Inches(0.9))
        self._para(tt, title, 32, INDIGO_DK, bold=True, first=True)
        self._rect(sl, Inches(0.97), Inches(1.52), Inches(2.0), Inches(0.055), INDIGO)
        ld = self._tbox(sl, Inches(0.97), Inches(1.72), Inches(11.6), Inches(0.9))
        self._para(ld, lede, 18, INK, first=True)
        body = self._tbox(sl, Inches(0.97), Inches(2.75), Inches(11.7), Inches(4.0))
        after_main = 11 if len(bullets) >= 8 else 13
        for i, b in enumerate(bullets):
            capt = b.startswith("~"); text = b[1:].strip() if capt else b
            if capt:
                self._para(body, text, 18, INDIGO_DK, bold=True, after=10, first=(i == 0))
            else:
                self._para(body, marker + text, 19, INK, after=after_main, first=(i == 0))
        ft = self._tbox(sl, Inches(0.97), Inches(6.98), Inches(9.0), Inches(0.4))
        self._para(ft, "Building With Claude", 11, GREY, first=True)
        pn = self._tbox(sl, Inches(11.4), Inches(6.98), Inches(1.2), Inches(0.4))
        self._para(pn, str(idx), 11, GREY, align=PP_ALIGN.RIGHT, first=True)
        self._notes(sl, s["notes"])
        self._estimate_overflow(title, lede, bullets, marker, after_main, idx)

    def _estimate_overflow(self, title, lede, bullets, marker, after_main, idx):
        avail = 6.92 - 2.75; h = 0.0
        for b in bullets:
            capt = b.startswith("~"); text = (b[1:].strip() if capt else marker + b)
            cpl = 80 if capt else 88; lh = (18 if capt else 19) * 1.20 / 72.0
            after = (10 if capt else after_main) / 72.0
            h += max(1, math.ceil(len(text) / cpl)) * lh + after
        msg = f"slide {idx:>2} '{title}': est {h:.2f}in / {avail:.2f}in" + (
            " <-- TIGHT" if h > avail else "")
        print("  " + msg)
        if h > avail:
            self.warnings.append(msg)

    def render(self, slides):
        for s in slides:
            if s["kind"] == "title":
                self.title_slide(s)
            elif s["kind"] == "divider":
                self.divider_slide(s)
            else:
                self.content_slide(s)


# ======================== entry point =====================================

def main(argv=None):
    ap = argparse.ArgumentParser(description="Build the kickoff deck from markdown.")
    ap.add_argument("src", nargs="?", default=str(DEFAULT_SRC),
                    help="kickoff-deck.md source (default: the handbook copy)")
    ap.add_argument("--out", default=str(DEFAULT_OUT), help="output .pptx path")
    args = ap.parse_args(argv)

    md = Path(args.src).read_text(encoding="utf-8")
    slides = parse_deck(md)

    nums = [s["num"] for s in slides]
    if nums != list(range(1, len(nums) + 1)):
        sys.exit(f"ERROR: slide numbers are not contiguous from 1: {nums}")
    if len(slides) != EXPECTED_SLIDES:
        sys.exit(f"ERROR: parsed {len(slides)} slides, expected {EXPECTED_SLIDES}")

    deck = Deck()
    deck.render(slides)
    out = Path(args.out)
    deck.prs.save(str(out))

    print(f"saved {out}")
    print(f"slides: {len(deck.prs.slides._sldIdLst)}")
    print(f"TIGHT slides: {len(deck.warnings)}")
    for w in deck.warnings:
        print("  !!", w)


if __name__ == "__main__":
    main()
